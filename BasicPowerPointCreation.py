# -*- coding: utf-8 -*-
"""
Created on Wed Oct 16 20:14:20 2024

@author: kayla.green
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os 
import pandas as pd
import numpy as np
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR

MasterLocation = r"C:\Users\kayla.green\Documents\Tools\Masters"

#Military Standard Format
Title_Font = "Calibri"
Title_size = 42
SubTitle_Font = "Calibri"
SubTitle_size = 28
Title_FontColor = ""

Section_Font = "Calibri"
Section_size = 42
SubSection_Font = "Calibri"
SubSection_size = 28
Title_FontColor = ""

TitleHeader_Font = "Calibri"
TitleHeader_size = 32
Section_Font = "Calibri"
Section_size = 14
TableHeader_Font = "Calibri"
TableHeader_size = 14
TableBody_Font = "Calibri"
TableBody_size = 12
Nominal_FontColor = ""

Banner_Font = "Calibri"
Banner_size = 14
Banner_FontColor = ""

Classification_Font = ""
Classification_size = 14
ClassSecret_color = ""
ClassConfidential_color = ""
ClassCUI_Unclass_color = ""

def GetSlideId(prs, slideName):
    i = 0
    while i < len(prs.slide_layouts):
        if slideName in prs.slide_layouts[i].name:
            return i
        else:
            i = i + 1
    return -1

def GetShapeId(slide, shapeName):
    i = 0
    while i < len(slide.shapes):
        if shapeName in slide.shapes[i].name:
            return i
        else:
            i = i + 1
            
    if shapeName == "Subtitle":
        i = 0
        while i < len(slide.shapes):
            if "Text Placeholder" in slide.shapes[i].name:
                return i
            else:
                i = i + 1
                
    return -1
    
def PPTX_Title(prs, title_str, subtitle_str):

    # Title slide
    slide_id = -1
    slide_id = GetSlideId(prs, "Title Slide")
    slide_layout = prs.slide_layouts[slide_id]

    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_str
    title.text_frame.paragraphs[0].font.name = Title_Font
    
    shape_id = GetShapeId(slide, "Subtitle")
    
    if shape_id != -1:
        if subtitle_str != "":
            subtitle = slide.shapes[shape_id]
            subtitle.text = subtitle_str         
            subtitle.text_frame.paragraphs[0].font.name = SubTitle_Font
            
    return prs

def PPTX_SectionBreak(prs, title_str, subtitle_str = ""):

    # Section slide
    slide_id = -1
    slide_id = GetSlideId(prs, "Section Header")
    slide_layout = prs.slide_layouts[slide_id]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_str
    title.text_frame.paragraphs[0].font.name = Section_Font

    shape_id = GetShapeId(slide, "Subtitle")
    
    if shape_id != -1:
        if subtitle_str != "":
            subtitle = slide.shapes[shape_id]
            subtitle.text = subtitle_str
            subtitle.text_frame.paragraphs[0].font.name = SubSection_Font

    return prs

def PPTX_ContentSlide(prs, title_str = "", information = "", bulletpoints = [], level = [], photoChart = [], bannertxt = ""):
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    if title_str == "" and information == "" and len(bulletpoints) == 0 and len(photoChart) == 0:
        slide_id = -1
        slide_id = GetSlideId(prs, "Blank")
        slide_layout = prs.slide_layouts[slide_id]
        slide = prs.slides.add_slide(slide_layout)
    
    elif information == "" and len(bulletpoints) == 0 and len(photoChart) == 0:
        slide_id = -1
        slide_id = GetSlideId(prs, "Title Only")
        slide_layout = prs.slide_layouts[slide_id]
        slide = prs.slides.add_slide(slide_layout)
        
    elif len(photoChart) == 0:
        slide_id = -1
        slide_id = GetSlideId(prs, "Title and Content")
        slide_layout = prs.slide_layouts[slide_id]
        slide = prs.slides.add_slide(slide_layout)
        
        shape_id = GetShapeId(slide, "Content")
        if shape_id != -1:
            if information != "":
                text_frame = slide.placeholders[shape_id]
                text_frame.text = information
            
            if len(bulletpoints) != 0:
                shapes = slide.shapes
                content_box = shapes.placeholders[shape_id]
                text_frame = content_box.text_frame
                
                i = 0
                while i < len(bulletpoints):
                    p = text_frame.add_paragraph()

                    if len(level) !=0:
                        p.level = level[i] 
                        p.text = str(bulletpoints[i])
                    else:
                        p.level = 0 
                        p.text = str(bulletpoints[i])
                        
                    i = i + 1
            
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            left = slide.placeholders[shape_id].left
            top = slide.placeholders[shape_id].top
            width = slide.placeholders[shape_id].width
            height = slide.placeholders[shape_id].height

            # Check if text box goes off the page
            if left + width > slide_width or top + height > slide_height:
                print(f"Text box in slide {slide.slide_id} goes off the page.")

                    
    else:
        slide_id = -1
        slide_id = GetSlideId(prs, "Picture")
        slide_layout = prs.slide_layouts[slide_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        
        i = 0
        while i < len(photoChart):
            print(type(photoChart[i]))
            size = photoChart[i].getbuffer().nbytes
            print(f'The size of the BytesIO object is {size} bytes.')

            image_width = 32206
            image_height = 32206
            
            # Convert image dimensions to inches
            image_width_inches = Inches(image_width / 96)  # Assuming 96 DPI
            image_height_inches = Inches(image_height / 96)
            
            # Add the plot image to the slide from the BytesIO object
            # Calculate the position to center the image
            left = (slide_width - image_width_inches) / 2
            top = (slide_height - image_height_inches) / 2
            
            slide.shapes.add_picture(photoChart[i], left, top, width=Inches(6), height=Inches(4))
            i = i + 1
            
        shape_id = GetShapeId(slide, "Content")
        if shape_id != -1:
            if information != "":
                text_frame = slide.placeholders[shape_id]
                text_frame.text = information
            
            if len(bulletpoints) != 0:
                shapes = slide.shapes
                content_box = shapes.placeholders[1]
                text_frame = content_box.text_frame
                
                i = 0
                while i < len(bulletpoints):
                    p = text_frame.add_paragraph()

                    if len(level) !=0:
                        p.level = level[i] 
                        p.text = bulletpoints[i]
                    else:
                        p.level = 0 
                        p.text = bulletpoints[i]
                        
                    i = i + 1
                    
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
 
    title = slide.shapes.title
    title.text = title_str
    title.text_frame.paragraphs[0].font.name = TitleHeader_Font
    
    if bannertxt != "":
        slideWidth = prs.slide_width/914400.0
        slideHeight = prs.slide_height/914400.0

        banner_box = slide.shapes.add_textbox(Inches(slideWidth/2-.5), Inches(slideHeight-1.5), Inches(1), Inches(1))
        banner_text = banner_box.text_frame
    
        p = banner_text.add_paragraph()
        p.text = bannertxt
        banner_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        fill = banner_box.fill
        fill.solid()
        main_color = RGBColor(0, 112, 192)  # Example RGB value for a theme color
        fill.fore_color.rgb = main_color
        
        banner_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        banner_box.vertical_anchor = MSO_ANCHOR.TOP
        
    return prs

def PPTX_Creation(Title_str, subtitle_str = "", PPT_Master = ""):

    # Create a presentation object
    if PPT_Master == "":
        prs = Presentation() 
    else:
        prs = Presentation(PPT_Master) 
    
    PPTX_Title(prs, Title_str, subtitle_str) 
    
    return prs

def PPTX_Save(prs, savepath, openaftercreation = True):
    prs.save(savepath)
    
    if openaftercreation == True:
        os.startfile(savepath)